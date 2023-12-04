from PIL import Image
from pd2ppt import df_to_table
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement

class Ppt:
    # Equivalent to a top-level class
    title_dimensions = [1.11, 6.43, 23, 1.5]
    subtitle1_dimensions = [4.47, 8.05, 16.46, 1.2]
    subtitle2_dimensions = [8.68, 12.49, 16.46, 1.55]
    privacy_dimensions = [5.38, 15.41, 23.95, 0.99]
    template = None

    header_dimensions = [2.36, 0.61, 8.81, 1.16]
    pri_color = (15, 15, 15)

    class Slide:
        def __init__(self, prs, layout, pri_color):
            self.parent = prs
            self.data = self.parent.slides.add_slide(layout)
            self.pri_color = pri_color
            pass

        def add_centered_picture(self, image):
            im = Image.open(image)

            left = (13.33 - (im.width / im.info['dpi'][0])) / 2
            top = ((6.7 - (im.height / im.info['dpi'][1])) / 2) + 0.8

            pic = self.data.shapes.add_picture(image, Inches(left), Inches(top))

            # Send image to back
            self.data.shapes[0]._element.addprevious(pic._element)
            return pic

        def add_df_to_table(self, df, size, columns_width=None, dimensions=None, headers_color=(7, 55, 99),
                            headers_height=None, data_color=(255, 255, 255), font_family='Arial', font_color=(5, 5, 5),
                            shorten_table=True):
            if dimensions is None:
                dimensions = [0.4, 1.09, 5.46, 12.5]

            l, t, h, w = dimensions
            if shorten_table:
                h = (h/size) * len(df)

            table = df_to_table(self.data, df, left=Inches(l), top=Inches(t), height=Inches(h), width=Inches(w)).table
            table.horz_banding = False

            self.set_table_dimensions(table, columns_width, len(df.columns), headers_height)

            rows = range(0, len(df)+1)
            columns = range(0, len(df.columns))

            for row in rows:
                for column in columns:
                    current_cell = table.cell(row, column)

                    self._set_cell_border(current_cell)
                    current_cell = self.set_cell_color(current_cell, row, headers_color, data_color)

                    current_cell, para = self.format_text(current_cell, font_family, font_color)

                    if row == 0:
                        para.font.size = Pt(9)
                        para.font.color.rgb = RGBColor(252, 252, 252)
                    else:
                        para.font.size = Pt(9)

            return table

        def add_header(self, text, dimensions):
            text_frame = self.add_textbox(text, dimensions, font_size=20, bold=True, halign='left', fit=(True, 20), rgb=self.pri_color)
            return text_frame

        def add_picture(self, image, dimensions,crop_top=None, crop_bottom=None):
            l, t, w, h = dimensions
            pic = self.data.shapes.add_picture(image, Inches(l), Inches(t), width=Inches(w), height=Inches(h))
            if crop_top is not None:
                pic.crop_top = crop_top
            if crop_bottom is not None:
                pic.crop_bottom = crop_bottom
            self.data.shapes[0]._element.addprevious(pic._element)
            return pic

        def add_rectangle(self, dimensions, rgb_fill, rgb_line, shadow=True):
            l, t, w, h = dimensions
            shape = self.data.shapes.add_shape(
                autoshape_type_id=MSO_SHAPE.RECTANGLE,
                left=Inches(l),
                top=Inches(t),
                width=Inches(w),
                height=Inches(h)
            )
            if shadow is False:
                shape.shadow.inherit = False
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(rgb_fill[0], rgb_fill[1], rgb_fill[2])
            line = shape.line
            line.color.rgb = RGBColor(rgb_line[0], rgb_line[1], rgb_line[2])

            return shape

        def add_table(self, rows, columns, dimensions, cells_df, columns_width=None, headers_color=(239, 239, 239),
                      data_color=(255, 255, 255), font_color=(0, 0, 0), font_family='Arial', headers_height=None):
            l, t, w, h = dimensions
            table = self.data.shapes.add_table(rows, columns, Inches(l), Inches(t), Inches(w), Inches(h)).table
            table.horz_banding = False

            self.set_table_dimensions(table, columns_width, columns, headers_height)

            rows = range(0, rows)
            columns = range(0, columns)

            for row in rows:
                for column in columns:
                    key = str(row) + str(column)
                    current_cell = table.cell(row, column)
                    self._set_cell_border(current_cell)

                    current_cell = self.set_cell_color(current_cell, row, headers_color, data_color)

                    if cells_df[key] is None:
                        continue

                    current_cell.text = cells_df[key]['text']

                    self.format_text(current_cell, font_family, font_color)

                    try:
                        if cells_df[key]['bold']:
                            current_cell.text_frame.paragraphs[0].font.bold = True
                    except KeyError:
                        pass

            return table

        def add_textbox(self, text, dimensions, cm=True, font_size=12, font_name='Arial', bold=False, italic=None,
                        rgb=(0, 0, 0), halign='center', valign='top', margin=None, fit=None):

            l, t, w, h = dimensions

            if cm:
                l = self.cm_to_inch_conversion(l)
                t = self.cm_to_inch_conversion(t)
                w = self.cm_to_inch_conversion(w)
                h = self.cm_to_inch_conversion(h)

            txtbox = self.data.shapes.add_textbox(left=Inches(l), top=Inches(t), width=Inches(w), height=Inches(h))
            text_frame = txtbox.text_frame
            text_frame.clear()
            text_frame = self.textbox(text_frame, text, font_size, font_name, bold, italic, rgb, halign, valign, margin)
            text_frame = self.auto_fit(text_frame, fit, font_size, bold)

            return text_frame

        def add_textbox_to_shape(self, existing_shape, text, font_size=12, font_name='Arial', bold=False, italic=None,
                                 rgb=None, halign='center', valign='top', margin=None, fit=None):
            if rgb is None:
                rgb = self.pri_color
            text_frame = existing_shape.text_frame
            text_frame = self.textbox(text_frame, text, font_size, font_name, bold, italic, rgb, halign, valign, margin)
            text_frame = self.auto_fit(text_frame, fit, font_size, bold)

            return text_frame

        @staticmethod
        def _set_cell_border(cell, border_color="000000", border_width='12700'):
            def sub_element(parent, tagname, **kwargs):
                element = OxmlElement(tagname)
                element.attrib.update(kwargs)
                parent.append(element)
                return element

            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for lines in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
                ln = sub_element(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
                solidFill = sub_element(ln, 'a:solidFill')
                srgbClr = sub_element(solidFill, 'a:srgbClr', val=border_color)
                prstDash = sub_element(ln, 'a:prstDash', val='solid')
                round_ = sub_element(ln, 'a:round')
                headEnd = sub_element(ln, 'a:headEnd', type='none', w='med', len='med')
                tailEnd = sub_element(ln, 'a:tailEnd', type='none', w='med', len='med')

        @staticmethod
        def auto_fit(text_frame, fit, font_size, bold):
            if fit is None:
                fit = (False, font_size)
            to_fit, maxsize = fit
            if isinstance(to_fit, bool) and isinstance(maxsize, (int, float)):
                if to_fit:
                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    if bold:
                        font_file = 'arialbd'
                    else:
                        font_file = 'arial'
                    try:
                        text_frame.fit_text(font_family='Arial', max_size=maxsize, bold=bold, font_file='./ktm_ppt/powerpoint-class/assets/fonts/{0}.ttf'.format(font_file))
                    except OSError:
                        print('Failed to fit text. \n Please include font folder with the following font types "arial.ttf, arialbd.ttf" \n under assets/fonts')
                    except TypeError:
                        print('Failed to fit text. \n Failed to break line as there are no spaces in text')
            else:
                raise ValueError('Parameters for fit must be of format (bool, int)')
            return text_frame

        @staticmethod
        def cm_to_inch_conversion(n):
            return n * (1 / 2.54)

        @staticmethod
        def format_text(cell, font_family, font_color):
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            para.font.name = font_family
            r, g, b = font_color
            para.font.color.rgb = RGBColor(r, g, b)

            return cell, para

        @staticmethod
        def reshape(shape, dimensions):
            l, t, w, h = dimensions
            shape.left = Inches(l)
            shape.top = Inches(t)
            shape.width = Inches(w)
            shape.height = Inches(h)

            return shape

        @staticmethod
        def set_table_dimensions(table, width_array, column_len, headers_height):
            if headers_height is not None:
                table.rows[0].height = Inches(headers_height)

            if width_array is not None:
                if len(width_array) == column_len:
                    for ind, _ in enumerate(width_array):
                        table.columns[ind].width = Inches(_)
                else:
                    raise ValueError("Width array must be equal to the number of columns")
            return table

        @staticmethod
        def set_cell_color(cell, row, headers_color, data_color):
            if row == 0:
                r, g, b = headers_color
            else:
                r, g, b = data_color

            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(r, g, b)
            return cell

        @staticmethod
        def set_individual_cell_color(cell, color):
            r, g, b = color
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(r, g, b)

        @staticmethod
        def textbox(text_frame, text, font_size, font_name, bold, italic, rgb, halign, valign, margin):
            # Vertical Alignment
            if valign == 'middle':
                text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            elif valign == 'bottom':
                text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.BOTTOM
            elif valign == 'top':
                text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

            p = text_frame.paragraphs[0]

            # Horizontal Alignment
            if halign == 'center':
                p.alignment = PP_ALIGN.CENTER
            elif halign == 'left':
                p.alignment = PP_ALIGN.LEFT
            elif halign == 'right':
                p.alignment = PP_ALIGN.RIGHT

            if margin:
                text_frame.margin_bottom = Inches(margin)
                text_frame.margin_top = Inches(margin)
                text_frame.margin_left = Inches(margin)
                text_frame.margin_right = Inches(margin)

            run = p.add_run()
            run.text = text
            font = run.font
            font.name = font_name
            font.size = Pt(font_size)
            font.bold = bold
            font.italic = italic
            R, G, B = rgb
            font.color.rgb = RGBColor(R, G, B)

            return text_frame

    def __init__(self, title, subtitle1=None, subtitle2=None, privacy=None, slides=None):
        if self.template is None:
            self.data = Presentation()
            self.titleLayout = self.data.slide_masters[0].slide_layouts[0]
            self.layout = self.data.slide_masters[0].slide_layouts[1]
        else:
            self.data = Presentation(self.template)
            self.titleLayout = self.data.slide_masters[0].slide_layouts[0]
            self.layout = self.data.slide_masters[0].slide_layouts[1]

        # Delete first & second slide
        try:
          self.delete_last_slide()
          self.delete_last_slide()
        except IndexError:
          print('No slides to delete')

        if slides is None:
            self.slides = []
            self.title_slide = self.add_title_slide(title, subtitle1, subtitle2, privacy)
        else:
            self.slides = slides

    def add_slide(self, title, master=0, layout=0):
        if layout:
            self.layout = self.data.slide_masters[master].slide_layouts[layout]

        new_slide = self.Slide(self.data, self.layout, self.pri_color)
        self.slides.append(new_slide)
        new_slide.add_header(title, self.header_dimensions)
        return new_slide

    def add_title_slide(self, title, subtitle1, subtitle2, privacy):
        title_slide = self.Slide(self.data, self.titleLayout, self.pri_color)
        self.slides.append(title_slide)
        title_slide.add_textbox(title, self.title_dimensions, font_size=33, bold=True, rgb=(63, 63, 63))

        if subtitle1 is not None:
            title_slide.add_textbox(subtitle1, self.subtitle1_dimensions, font_size=18, rgb=(150, 146, 146))

        if subtitle2 is not None:
            title_slide.add_textbox(subtitle2, self.subtitle2_dimensions, font_size=18, rgb=(127, 127, 127))

        if privacy is not None:
            privacy_frame = title_slide.add_rectangle(self.privacy_dimensions, rgb_fill=[255, 255, 0], rgb_line=[0, 0, 0])
            title_slide.add_textbox_to_shape(privacy_frame, privacy, font_size=14, italic=True, rgb=[0, 0, 0])
        return title_slide

    def delete_last_slide(self):
        # Delete Second slide
        rId = self.data.slides._sldIdLst[-1].rId
        self.data.part.drop_rel(rId)
        del self.data.slides._sldIdLst[-1]

    def save(self, directory):
        self.data.save(directory)

    def set_layout(self, n):
        self.layout = self.data.slide_masters[0].slide_layouts[n]

# text_frame = self.add_textbox(title, dimensions, font_size=36, bold=True, rgb=[63, 63, 63])
